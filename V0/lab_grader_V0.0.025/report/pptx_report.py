"""
Requirement 6 (PowerPoint Export) — python-pptx로 템플릿(.pptx)의 {{PLACEHOLDER}} 텍스트를
실제 값으로 치환. 반드시 run.text 레벨에서만 바꿔서 디자이너가 잡아둔 폰트/색/굵기(run 서식)를
그대로 보존한다 — paragraph.text = ... 로 통째로 바꾸면 그 문단의 서식이 깨지므로 절대 쓰지 않는다.
"""
import os
import re
from pptx import Presentation

_PLACEHOLDER_RE = re.compile(r"\{\{\s*([A-Z0-9_]+)\s*\}\}")

UNREACHABLE_LABEL = "접속 불가"


def build_placeholder_map(dataset, device_name=None):
    """
    dataset: {device_name: metrics_dict} (report.textfsm_parser.build_report_dataset 결과)
    device_name이 주어지면 그 장비 metrics만, 없으면 dataset에 장비가 1개일 때만 사용.
    {{CPU_UTIL}}, {{MEMORY_UTIL}}, {{MODEL}}, {{SERIAL_NUMBER}}, {{IMAGE}}, {{UPTIME}},
    {{DEVICE_NAME}} 키로 매핑.
    """
    if device_name is None:
        if len(dataset) != 1:
            raise ValueError("device_name을 지정하거나 dataset에 장비가 1개여야 합니다.")
        device_name = next(iter(dataset))
    metrics = dataset.get(device_name, {})

    if metrics.get("unreachable"):
        keys = ["MODEL", "SERIAL_NUMBER", "IMAGE", "UPTIME", "CPU_UTIL", "MEMORY_UTIL", "LOAD_AVERAGE"]
        mapping = {k: UNREACHABLE_LABEL for k in keys}
    else:
        mapping = {
            "MODEL": str(metrics.get("model", "")),
            "SERIAL_NUMBER": str(metrics.get("serial_number", "")),
            "IMAGE": str(metrics.get("image", "")),
            "UPTIME": str(metrics.get("uptime", "")),
            "CPU_UTIL": _fmt_percent(metrics.get("cpu_util_percent")),
            "MEMORY_UTIL": _fmt_percent(metrics.get("memory_used_percent")),
            "LOAD_AVERAGE": str(metrics.get("load_average_1m", "")),
        }
    mapping["DEVICE_NAME"] = device_name
    return mapping


def _fmt_percent(value):
    return f"{value}%" if value not in (None, "") else ""


def _replace_in_runs(runs, mapping):
    """
    {{KEY}}가 여러 run에 걸쳐 쪼개져 있을 수도 있으므로, 문단 전체 텍스트를 합쳐서
    치환 대상 여부만 먼저 판단하고, 실제 교체는 "플레이스홀더 전체가 한 run 안에 있는
    표준적인 경우"에 한해 run.text만 바꾼다(서식 보존). 쪼개진 경우는 손대지 않고 원본 유지.
    """
    for run in runs:
        if not run.text or "{{" not in run.text:
            continue

        def repl(match):
            key = match.group(1)
            return mapping.get(key, match.group(0))

        new_text = _PLACEHOLDER_RE.sub(repl, run.text)
        if new_text != run.text:
            run.text = new_text


def apply_placeholders_to_pptx(input_path, output_path, mapping):
    prs = Presentation(input_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                _replace_in_runs(paragraph.runs, mapping)
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return output_path


def build_blank_template(output_path):
    """템플릿이 아직 없는 최초 실행을 위한 최소 {{PLACEHOLDER}} 템플릿 생성."""
    from pptx.util import Inches, Pt

    prs = Presentation()
    layout = prs.slide_layouts[5]  # 제목만 있는 레이아웃
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "{{DEVICE_NAME}} 점검 결과"

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.5), Inches(4))
    tf = box.text_frame
    lines = [
        "모델: {{MODEL}}", "시리얼: {{SERIAL_NUMBER}}", "소프트웨어: {{IMAGE}}",
        "가동시간: {{UPTIME}}", "CPU 사용률: {{CPU_UTIL}}", "메모리 사용률: {{MEMORY_UTIL}}",
    ]
    tf.text = lines[0]
    tf.paragraphs[0].runs[0].font.size = Pt(18)
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.runs[0].font.size = Pt(18)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    tpl = "AutoCheck_pptx_template_sample.pptx"
    out = "AutoCheck_report_sample.pptx"
    build_blank_template(tpl)
    mapping = build_placeholder_map({
        "Core1": {"model": "DCS-7050SX3-48YC8", "serial_number": "SN1", "image": "4.28.0F",
                  "uptime": "3 weeks", "cpu_util_percent": 16.0, "memory_used_percent": 62.5,
                  "load_average_1m": "0.10"},
    })
    apply_placeholders_to_pptx(tpl, out, mapping)
    print("작성됨:", out)
